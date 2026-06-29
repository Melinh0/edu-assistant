import requests
import os
import platform
import streamlit as st
from PyPDF2 import PdfReader
from io import BytesIO
import base64
import fitz
import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from docx import Document as DocxDocument

OCR_SPACE_API_KEY = os.getenv("OCR_SPACE_API_KEY")
OCR_URL = "https://api.ocr.space/parse/image"
OLLAMA_API_KEY = os.getenv("OLLAMA_API_KEY")
OLLAMA_HOST = os.getenv("OLLAMA_HOST")

def setup_tesseract():
    if platform.system() == "Windows":
        tesseract_paths = [
            r"C:\Program Files\Tesseract-OCR\tesseract.exe",
            r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
            r"C:\ProgramData\chocolatey\bin\tesseract.exe"
        ]
        tessdata_dirs = [
            r"C:\Program Files\Tesseract-OCR\tessdata",
            r"C:\Program Files (x86)\Tesseract-OCR\tessdata"
        ]
    else:
        tesseract_paths = ["/usr/bin/tesseract", "/usr/local/bin/tesseract"]
        tessdata_dirs = [
            "/usr/share/tesseract-ocr/4.00/tessdata",
            "/usr/share/tesseract-ocr/tessdata",
            "/usr/local/share/tesseract-ocr/tessdata"
        ]

    for path in tesseract_paths:
        if os.path.exists(path):
            import pytesseract
            pytesseract.pytesseract.tesseract_cmd = path
            break

    tessdata_dir = None
    if "TESSDATA_PREFIX" in os.environ and os.path.exists(os.environ["TESSDATA_PREFIX"]):
        tessdata_dir = os.environ["TESSDATA_PREFIX"]
    else:
        for d in tessdata_dirs:
            if os.path.exists(d):
                tessdata_dir = d
                os.environ["TESSDATA_PREFIX"] = d
                break

    if tessdata_dir:
        lang_file = os.path.join(tessdata_dir, "por.traineddata")
        if not os.path.exists(lang_file):
            try:
                st.info("Baixando idioma português para o Tesseract...")
                url = "https://github.com/tesseract-ocr/tessdata/raw/main/por.traineddata"
                response = requests.get(url, timeout=30)
                if response.status_code == 200:
                    with open(lang_file, "wb") as f:
                        f.write(response.content)
                    st.success("Idioma português baixado com sucesso.")
                else:
                    st.warning("Não foi possível baixar o idioma português. O OCR pode falhar.")
            except Exception as e:
                st.warning(f"Erro ao baixar idioma: {e}")

try:
    import pytesseract
    setup_tesseract()
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False
    st.warning("Tesseract não instalado. Instale com: pip install pytesseract")

def preprocess_image(pil_img):
    img = np.array(pil_img.convert('RGB'))
    h, w = img.shape[:2]
    if w < 1500:
        scale = max(2000 / w, 3.0)
        new_w = int(w * scale)
        new_h = int(h * scale)
        img = cv2.resize(img, (new_w, new_h), interpolation=cv2.INTER_CUBIC)
    gray = cv2.cvtColor(img, cv2.COLOR_RGB2GRAY)
    gray = cv2.fastNlMeansDenoising(gray, None, 15, 7, 21)
    clahe = cv2.createCLAHE(clipLimit=4.0, tileGridSize=(8,8))
    gray = clahe.apply(gray)
    gray = cv2.bilateralFilter(gray, 9, 75, 75)
    kernel = np.array([[-1,-1,-1], [-1,9,-1], [-1,-1,-1]])
    gray = cv2.filter2D(gray, -1, kernel)
    thresh = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 21, 4)
    kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (3,3))
    thresh = cv2.morphologyEx(thresh, cv2.MORPH_CLOSE, kernel, iterations=1)
    return Image.fromarray(thresh)

def postprocess_text(text):
    replacements = {
        '|': 'l', '1': 'l', '0': 'O', '5': 'S',
        'é': 'e', 'ã': 'a', 'ç': 'c', 'á': 'a',
        'í': 'i', 'ó': 'o', 'ú': 'u', 'â': 'a',
        'ê': 'e', 'ô': 'o', '"': '', '`': "'",
        '…': '...', '—': '-', '–': '-'
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    lines = text.split('\n')
    clean_lines = []
    for line in lines:
        line = line.strip()
        if line:
            if line.endswith('-'):
                line = line[:-1]
            if len(line) > 1:
                clean_lines.append(line)
    return '\n'.join(clean_lines)

def refine_text_with_llm(raw_text, model, cache_buster=None):
    if not raw_text or not model:
        return raw_text
    try:
        from utils.llm import call_ollama
        prompt = f"""O texto abaixo foi extraído de um PDF escaneado por OCR e está muito bagunçado, com palavras quebradas, caracteres trocados e estrutura confusa. 
Sua tarefa é **reorganizar, corrigir e estruturar** este texto para que fique legível e organizado, mantendo o conteúdo original.
- Corrija palavras com erros comuns (ex: 'crianqa' -> 'criança', 'velocidade inlcial' -> 'velocidade inicial').
- Organize o texto em parágrafos ou tópicos coerentes.
- Se houver questões, mantenha a numeração e as alternativas.
- Se houver figuras, mantenha a descrição textual.
- Preserve o significado original. Não invente informações.
- Responda APENAS com o texto corrigido e organizado.

Texto bruto do OCR:
{raw_text}"""
        refined = call_ollama(prompt, model, cache_buster=cache_buster)
        return refined if refined else raw_text
    except Exception as e:
        st.warning(f"Falha ao refinar texto com LLM: {e}")
        return raw_text

def extract_text_from_file(file_bytes, filename, model=None, cache_buster=None):
    ext = filename.split('.')[-1].lower()
    if ext == 'pdf':
        return extract_text_from_pdf(file_bytes, model, cache_buster=cache_buster)
    elif ext in ['png', 'jpg', 'jpeg']:
        return extract_text_from_image(file_bytes)
    elif ext == 'pptx':
        return extract_text_from_pptx(file_bytes)
    elif ext == 'docx':
        return extract_text_from_docx(file_bytes)
    elif ext == 'txt':
        return file_bytes.decode('utf-8', errors='ignore')
    else:
        st.warning(f"Formato {ext} não suportado para extração.")
        return ""

def extract_text_from_pdf(pdf_bytes, model=None, cache_buster=None):
    raw_text = ""
    try:
        reader = PdfReader(BytesIO(pdf_bytes))
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                raw_text += page_text + "\n"
        if raw_text.strip():
            raw_text = postprocess_text(raw_text)
    except:
        pass

    if not raw_text.strip():
        try:
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            for page in doc:
                page_text = page.get_text()
                if page_text:
                    raw_text += page_text + "\n"
            doc.close()
            if raw_text.strip():
                raw_text = postprocess_text(raw_text)
        except:
            pass

    if not raw_text.strip() and TESSERACT_AVAILABLE:
        st.info("Usando Tesseract OCR para extração...")
        try:
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            all_text = ""
            for page_num, page in enumerate(doc, start=1):
                pix = page.get_pixmap(dpi=300)
                img_bytes = pix.tobytes("png")
                img = Image.open(BytesIO(img_bytes))
                processed = preprocess_image(img)
                config = '--psm 6 --oem 3 -l por+eng'
                page_text = pytesseract.image_to_string(processed, config=config)
                if page_text.strip():
                    all_text += f"\n--- Página {page_num} ---\n{page_text}\n"
            doc.close()
            if all_text.strip():
                raw_text = postprocess_text(all_text)
        except Exception as e:
            st.error(f"Tesseract falhou: {str(e)}")

    if not raw_text.strip():
        raw_text = extract_text_with_ocr(pdf_bytes)
        if raw_text.strip():
            raw_text = postprocess_text(raw_text)

    if raw_text.strip():
        if model:
            st.info("Refinando texto extraído com o LLM...")
            return refine_text_with_llm(raw_text, model, cache_buster=cache_buster)
        else:
            return raw_text

    if model and is_vision_model(model):
        st.info("Tentando extrair texto com visão do Ollama...")
        vision_text = extract_text_with_vision(pdf_bytes, model)
        if vision_text.strip():
            return refine_text_with_llm(vision_text, model, cache_buster=cache_buster) if model else vision_text
    else:
        st.warning("Não foi possível extrair texto. Selecione um modelo com suporte a visão (ex: gemma4:31b).")
        return ""

def extract_text_from_image(file_bytes):
    if TESSERACT_AVAILABLE:
        img = Image.open(BytesIO(file_bytes))
        processed = preprocess_image(img)
        config = '--psm 6 --oem 3 -l por+eng'
        text = pytesseract.image_to_string(processed, config=config)
        return postprocess_text(text)
    return ""

def extract_text_from_pptx(file_bytes):
    try:
        prs = Presentation(BytesIO(file_bytes))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                if hasattr(shape, "notes_slide") and shape.notes_slide:
                    notes = shape.notes_slide.notes_text_frame.text
                    if notes:
                        text += notes + "\n"
        return postprocess_text(text)
    except Exception as e:
        st.error(f"Erro ao extrair texto do PowerPoint: {str(e)}")
        return ""

def extract_text_from_docx(file_bytes):
    try:
        doc = DocxDocument(BytesIO(file_bytes))
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return postprocess_text(text)
    except Exception as e:
        st.error(f"Erro ao extrair texto do DOCX: {str(e)}")
        return ""

def extract_text_with_ocr(pdf_bytes):
    if not OCR_SPACE_API_KEY:
        st.error("Chave OCR Space não configurada.")
        return ""
    if len(pdf_bytes) > 1.5 * 1024 * 1024:
        st.warning("Arquivo > 1.5 MB. OCR Space (plano gratuito) não suporta. Usando fallback...")
        return ""
    files = {'file': ('document.pdf', pdf_bytes, 'application/pdf')}
    data = {
        'apikey': OCR_SPACE_API_KEY,
        'language': 'por',
        'isOverlayRequired': False,
        'filetype': 'PDF'
    }
    try:
        response = requests.post(OCR_URL, files=files, data=data, timeout=60)
        if response.status_code == 200:
            result = response.json()
            if result.get('OCRExitCode') == 1:
                return result['ParsedResults'][0]['ParsedText']
            else:
                st.error(f"Erro no OCR: {result.get('ErrorMessage', 'Desconhecido')}")
                return ""
        else:
            st.error(f"Falha na requisição OCR: {response.status_code}")
            st.error(f"Detalhes: {response.text}")
            return ""
    except Exception as e:
        st.error(f"Exceção no OCR: {str(e)}")
        return ""

def is_vision_model(model):
    vision_models = os.getenv("OLLAMA_VISION_MODELS", "gemma4:31b").split(',')
    vision_models = [m.strip() for m in vision_models if m.strip()]
    return model in vision_models

def extract_text_with_vision(pdf_bytes, model):
    if not OLLAMA_API_KEY:
        st.error("Chave Ollama não configurada.")
        return ""
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        all_text = ""
        max_pages = 10
        for i, page in enumerate(doc):
            if i >= max_pages:
                st.info(f"Limite de {max_pages} páginas para visão. As demais não foram processadas.")
                break
            pix = page.get_pixmap(dpi=200)
            img_bytes = pix.tobytes("png")
            img_base64 = base64.b64encode(img_bytes).decode('utf-8')
            prompt = "Extraia todo o texto desta imagem. Responda apenas com o texto extraído, sem comentários adicionais."
            payload = {
                "model": model,
                "prompt": prompt,
                "images": [img_base64],
                "stream": False,
                "options": {"num_predict": 4000}
            }
            headers = {
                "Authorization": f"Bearer {OLLAMA_API_KEY}",
                "Content-Type": "application/json"
            }
            url = f"{OLLAMA_HOST}/api/generate"
            response = requests.post(url, json=payload, headers=headers, timeout=120)
            if response.status_code == 200:
                data = response.json()
                page_text = data.get('response', '')
                all_text += f"\n--- Página {i+1} ---\n{page_text}\n"
            else:
                st.error(f"Erro na chamada Ollama (página {i+1}): {response.status_code} - {response.text}")
        doc.close()
        return all_text.strip()
    except Exception as e:
        st.error(f"Erro ao extrair com visão: {str(e)}")
        return ""